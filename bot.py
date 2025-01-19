from telegram import Update
from telegram.ext import Application, CommandHandler, MessageHandler, filters, CallbackContext, ConversationHandler
from PIL import Image, ImageDraw
import pandas as pd
import os
from PyPDF2 import PdfReader
from pdf2image import convert_from_path
from fpdf import FPDF
from docx import Document
from docx import Document as DocxDocument
from pptx import Presentation
from io import BytesIO

# Define conversation states
SELECT_FORMAT, CONVERT_FILE = range(2)

# Define the start function
async def start(update: Update, context: CallbackContext) -> int:
    await update.message.reply_text('Hello! Please choose the format you want to convert to (e.g., jpeg, pdf, txt, docx, xlsx, ppt):')
    return SELECT_FORMAT

# Handle PowerPoint to PDF conversion
def handle_pptx_to_pdf(file_path):
    try:
        prs = Presentation(file_path)
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)

        for slide in prs.slides:
            slide_text = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    slide_text += shape.text + "\n"

            pdf.add_page()
            pdf.set_font("Arial", size=12)
            pdf.multi_cell(0, 10, slide_text)

        new_file_path = os.path.splitext(file_path)[0] + '.pdf'
        pdf.output(new_file_path)
        return new_file_path
    except Exception as e:
        print(f"Error converting PPTX to PDF: {e}")
        return None

# Handle PowerPoint to PowerPoint conversion
def handle_pptx_to_ppt(file_path):
    try:
        # Simply rename or copy the file with a new extension
        new_file_path = os.path.splitext(file_path)[0] + '.ppt'
        os.rename(file_path, new_file_path)
        return new_file_path
    except Exception as e:
        print(f"Error converting PPTX to PPT: {e}")
        return None

# Handle image file conversion (PNG to JPEG example)
def handle_image(file_path, new_format='JPEG'):
    try:
        with Image.open(file_path) as img:
            new_file_path = os.path.splitext(file_path)[0] + f'.{new_format.lower()}'
            img.save(new_file_path, new_format)
        return new_file_path
    except Exception as e:
        print(f"Error converting image: {e}")
        return None

# Handle image to PDF conversion
def image_to_pdf(file_path):
    try:
        with Image.open(file_path) as img:
            new_file_path = os.path.splitext(file_path)[0] + '.pdf'
            img.save(new_file_path, 'PDF', resolution=100.0)
        return new_file_path
    except Exception as e:
        print(f"Error converting image to PDF: {e}")
        return None

# Handle PDF to text conversion
def handle_pdf_to_text(file_path):
    try:
        pdf_reader = PdfReader(file_path)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text()
        new_file_path = os.path.splitext(file_path)[0] + '.txt'
        with open(new_file_path, 'w') as text_file:
            text_file.write(text)
        return new_file_path
    except Exception as e:
        print(f"Error converting PDF to text: {e}")
        return None

# Handle PDF to DOCX conversion
def handle_pdf_to_docx(file_path):
    try:
        pdf_reader = PdfReader(file_path)
        new_file_path = os.path.splitext(file_path)[0] + '.docx'
        doc = DocxDocument()
        for page in pdf_reader.pages:
            text = page.extract_text()
            doc.add_paragraph(text)
        doc.save(new_file_path)
        return new_file_path
    except Exception as e:
        print(f"Error converting PDF to DOCX: {e}")
        return None

# Handle PDF to JPEG conversion
def handle_pdf_to_jpeg(file_path):
    try:
        images = convert_from_path(file_path)
        new_file_paths = []
        for i, image in enumerate(images):
            new_file_path = os.path.splitext(file_path)[0] + f'_{i+1}.jpeg'
            image.save(new_file_path, 'JPEG')
            new_file_paths.append(new_file_path)
        return new_file_paths
    except Exception as e:
        print(f"Error converting PDF to JPEG: {e}")
        return None

# Handle DOCX to text conversion
def handle_docx_to_text(file_path):
    try:
        document = Document(file_path)
        text = ""
        for paragraph in document.paragraphs:
            text += paragraph.text + "\n"
        new_file_path = os.path.splitext(file_path)[0] + '.txt'
        with open(new_file_path, 'w') as text_file:
            text_file.write(text)
        return new_file_path
    except Exception as e:
        print(f"Error converting DOCX to text: {e}")
        return None

# Handle DOCX to JPEG conversion
def handle_docx_to_jpeg(file_path):
    try:
        document = Document(file_path)
        text = ""
        for paragraph in document.paragraphs:
            text += paragraph.text + "\n"
        img = Image.new('RGB', (1000, 1000), color=(255, 255, 255))
        d = ImageDraw.Draw(img)
        d.text((10, 10), text, fill=(0, 0, 0))
        new_file_path = os.path.splitext(file_path)[0] + '.jpeg'
        img.save(new_file_path)
        return new_file_path
    except Exception as e:
        print(f"Error converting DOCX to JPEG: {e}")
        return None

# Handle CSV to Excel conversion
def handle_csv(file_path):
    try:
        df = pd.read_csv(file_path)
        new_file_path = os.path.splitext(file_path)[0] + '.xlsx'
        df.to_excel(new_file_path, index=False)
        return new_file_path
    except Exception as e:
        print(f"Error converting CSV to Excel: {e}")
        return None

# Handle text file to PDF conversion
def text_to_pdf(file_path):
    try:
        pdf = FPDF()
        pdf.add_page()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.set_font("Arial", size=12)

        with open(file_path, 'r') as file:
            for line in file:
                pdf.cell(200, 10, txt=line, ln=True)

        new_file_path = os.path.splitext(file_path)[0] + '.pdf'
        pdf.output(new_file_path)
        return new_file_path
    except Exception as e:
        print(f"Error converting text to PDF: {e}")
        return None

# Handle DOCX file to PDF conversion
def docx_to_pdf(file_path):
    try:
        document = Document(file_path)
        pdf = FPDF()
        pdf.add_page()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.set_font("Arial", size=12)

        for paragraph in document.paragraphs:
            pdf.cell(200, 10, txt=paragraph.text, ln=True)

        new_file_path = os.path.splitext(file_path)[0] + '.pdf'
        pdf.output(new_file_path)
        return new_file_path
    except Exception as e:
        print(f"Error converting DOCX to PDF: {e}")
        return None

# Store user-selected format and proceed with conversion
async def select_format(update: Update, context: CallbackContext) -> int:
    user_format = update.message.text.lower()
    context.user_data['format'] = user_format
    await update.message.reply_text("Great! Now send me the file you want to convert.")
    return CONVERT_FILE

# Define the file handler function
async def handle_file(update: Update, context: CallbackContext) -> int:
    document = update.message.document
    file = await document.get_file()
    file_path = file.file_path
    file_name = document.file_id + "." + document.file_name.split('.')[-1]
    await file.download_to_drive(file_name)
    print(f"Downloaded file to: {file_name}")

    # Get user-selected format
    user_format = context.user_data.get('format')


    # Determine the file type and process accordingly
    if document.mime_type.startswith('image/'):  # Images
        if user_format == 'pdf':
            new_file_path = image_to_pdf(file_name)
        else:
            new_file_path = handle_image(file_name, user_format.upper())
    elif document.mime_type == 'application/pdf':  # PDFs
        if user_format == 'txt':
            new_file_path = handle_pdf_to_text(file_name)
        elif user_format == 'docx':
            new_file_path = handle_pdf_to_docx(file_name)
        elif user_format == 'jpeg':
            new_file_path = handle_pdf_to_jpeg(file_name)
        else:
            await update.message.reply_text("Unsupported file format for PDF files.")
            return ConversationHandler.END
    elif document.mime_type == 'text/plain':  # Plain text files
        if user_format == 'pdf':
            new_file_path = text_to_pdf(file_name)
        else:
            await update.message.reply_text("Unsupported file format for text files.")
            return ConversationHandler.END
    elif document.mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':  # Word documents
        if user_format == 'pdf':
            new_file_path = docx_to_pdf(file_name)
        elif user_format == 'txt':
            new_file_path = handle_docx_to_text(file_name)
        elif user_format == 'jpeg':
            new_file_path = handle_docx_to_jpeg(file_name)
        else:
            await update.message.reply_text("Unsupported file format for Word documents.")
            return ConversationHandler.END
    elif document.mime_type == 'text/csv':  # CSV files
        if user_format == 'xlsx':
            new_file_path = handle_csv(file_name)
        else:
            await update.message.reply_text("Unsupported file format for CSV files.")
            return ConversationHandler.END
    elif document.mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':  # PowerPoint presentations
        if user_format == 'pdf':
            new_file_path = handle_pptx_to_pdf(file_name)
        elif user_format == 'ppt':
            new_file_path = handle_pptx_to_ppt(file_name)
        else:
            await update.message.reply_text("Unsupported file format for PowerPoint files.")
            return ConversationHandler.END
    else:
        await update.message.reply_text("Unsupported file format.")
        return ConversationHandler.END

    # Send the converted file back to the user
    if new_file_path:
        if isinstance(new_file_path, list):  # Handle multiple files case (e.g., PDF to JPEG)
            for path in new_file_path:
                await update.message.reply_document(document=open(path, 'rb'))
                os.remove(path)
        else:
            await update.message.reply_document(document=open(new_file_path, 'rb'))
            os.remove(new_file_path)
        os.remove(file_name)
    else:
        await update.message.reply_text("There was an error converting the file.")

    return ConversationHandler.END

# Main function to run the bot
def main():
    # Replace 'YOUR_API_TOKEN' with the token you received from BotFather
    application = Application.builder().token('7565069768:AAGfUUlq3ptwVFitSQEjmr1JqaP_mQz1hbc').build()
    
    # Define conversation handler with states
    conv_handler = ConversationHandler(
        entry_points=[CommandHandler('start', start)],
        states={
            SELECT_FORMAT: [MessageHandler(filters.TEXT & ~filters.COMMAND, select_format)],
            CONVERT_FILE: [MessageHandler(filters.Document.ALL, handle_file)],
        },
        fallbacks=[],
    )

    # Add conversation handler to the application
    application.add_handler(conv_handler)

    # Start the bot
    application.run_polling()

if __name__ == '__main__':
    main()
